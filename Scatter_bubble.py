import plotly.graph_objects as go
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm, Inches
import plotly.express as px

def create_scatter_bubble_chart(details,base=None):
    try:
        prs = Presentation('./Final - Presentation.pptx')
    except exc.PackageNotFoundError as err:
        prs = Presentation(base)
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # print(px.data.iris().tail())
    x = details['scatter_bubble_data']['x'] 
    y = details['scatter_bubble_data']['y']
    z = details['scatter_bubble_data']['z']
    if len(x)!=len(y)!=len(z):
        return False, "X and Y do not have same length. Please enter the values again."
    fig = px.scatter( x=x, y=y, color=z,marginal_y="violin",
            marginal_x="box", trendline="ols", template="simple_white")
    
    fig.write_image('scatter_bubble_data2.jpg')
    shapes = prs.slides[0].shapes
    picture = shapes.add_picture('scatter_bubble_data2.jpg', Inches(0.1), Inches(0.1))

    prs.save("sample_presentation.pptx")


details = {
    "create_scatter_bubble_chart":True,
    "scatter_bubble_data":{
        "x":[2,4,6,8],
        "y":[4,6,8,20],
        "z":['red','green','red','green']
    }
}
# scatter_bubble(details)