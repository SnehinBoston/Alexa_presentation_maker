from pptx import Presentation,exc
import pandas as pd
from pptx.util import Inches,Pt
import matplotlib.pyplot as plt
import numpy as np
import matplotlib.colors as mcolors
import matplotlib.cm as cm

def negat(x):
    if int(x)<0:
        return True
    else:
        return False

def create_mekko_chart(details, base=None):
    
    zee = list(map(negat,details['mekko_data']['z']))
    if len(details['mekko_data']['x']) != len(details['mekko_data']['y'])+1 or len(details['mekko_data']['y']) != len(details['mekko_data']['z']):
        return (False, "Couldn't create the Mekko chart due to unequal lengths of x,y and z")
    elif True not in zee:
        return (False, "Couldn't create the Mekko chart due to negative value in z")
    
    x = np.array(details['mekko_data']['x'])
    y = np.array(details['mekko_data']['y'])
    z = np.array(details['mekko_data']['z'])
    
    print("x: ",len(x))
    print("y: ",len(y))
    print("z: ",len(z))

    bin_width = np.diff(x)
    fig, ax = plt.subplots()

    cmap = plt.cm.RdYlBu
    norm = mcolors.TwoSlopeNorm(vmin=z.min(), vcenter=0., vmax=z.max())
    sm = cm.ScalarMappable(cmap=cmap, norm=norm)
    sm.set_array([])

    bar_plot = ax.bar(x[:-1], y,  color=cmap(norm(z)),  width=bin_width,  align='edge')

    cbar = fig.colorbar(sm)

    ax.set_ylabel("y")

    fig1 = plt.gcf()
    fig1.savefig('mekko.png', dpi=100)

    try:
        prs = Presentation('./Final - Presentation.pptx')
    except exc.PackageNotFoundError as err:
        prs = Presentation(base)
    
    slide = prs.slides.add_slide(prs.slide_layouts[10])
    shapes = slide.shapes
    left = Inches(8)
    top =  Inches(1)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    tf.text = "Mekko chart"
    # title = shapes.title
    # title.text = "Mekko chart"
    left = Inches(3)
    top = Inches(2)
    pic = slide.shapes.add_picture("mekko.png", left, top,width = Inches(8),height = Inches(4.0))
    prs.save('./Final - Presentation.pptx')

    return True,"Created the Mekko chart"


details = {
    "create_mekko_chart":True,
    "mekko_data":{
        "x":[0, 0.4512135, 0.760715 , 0.775948 , 0.977063 , 1.170482 ,1.229812 , 1.3009845, 1.347207 , 1.4155705, 1.928897 ],
        "y":[0.40048296,  1.11131896,  0.30525134,  3.86793415, 21.80974083, 11.88354534, 13.84599687,  9.7484865 ,  9.5418679 , 22.39983675],
        "z":[0, 0.4512135, -0.760715 , 0.775948 , 0.977063 , 1.170482 ,1.229812 , 1.3009845, 1.347207 , -1.4155705]
    }
}

# print(create_mekko_chart(details))

req = { 
    "req":'create_chart',
    "details": {
        "create_mekko_chart":True,
        "mekko_data":{
            "x":[40.0,20.0,120.0,124.0,240.0,240.0,114.0,532.0,235.0],
            "y":[40.0,20.0,120.0,124.0,240.0,240.0,114.0,532.0,235.0],
            "z":[40.0,20.0,120.0,124.0,-240.0,240.0,114.0,532.0,235.0]
        }
    }
}