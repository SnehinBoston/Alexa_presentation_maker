from pptx import Presentation,exc
import pandas as pd
import plotly.express as px
from pptx.util import Inches,Pt
from datetime import datetime

def validate(date_text):
    try:
        if date_text != datetime.strptime(date_text, "%Y-%m-%d").strftime('%Y-%m-%d'):
            raise ValueError
        return True
    except ValueError:
        return False

def create_gantt_chart(details, base=None):
    # Need to write in some checks
    # Verify that Task are alphabets, Start & Finish are dates, Completion_pct is percentage.
    print("Hello")
    if len(details['gantt_data']['Task']) != len(details['gantt_data']['Completion_pct']):
        return (False, "Couldn't create the Gantt chart due to invalid lengths of Task and Completion percentage lists")
    elif len(details['gantt_data']['Start']) != len(details['gantt_data']['Finish']):
        return (False, "Couldn't create the Gantt chart due to invalid lengths of Start and Finish lists")
    
    Start_val = list(map(validate, details['gantt_data']['Start']))
    Finish_val = list(map(validate, details['gantt_data']['Finish']))
    
    if False in (Start_val):
        return (False, "Couldn't create the Gantt chart due to invalid formats of Start list")
    elif False in Finish_val:
        return (False, "Couldn't create the Gantt chart due to invalid formats of Finish list")

    df = pd.DataFrame(dict(Task=details['gantt_data']['Task'], Start=details['gantt_data']['Start'], Finish=details['gantt_data']['Finish'], Completion_pct=details['gantt_data']['Completion_pct']))
    fig = px.timeline(df, x_start="Start", x_end="Finish", y="Task", color="Completion_pct")
    fig.update_yaxes(autorange="reversed")
    fig.write_image("fig1.png")

    try:
        prs = Presentation('./Final - Presentation.pptx')
    except exc.PackageNotFoundError as err:
        prs = Presentation(base)
    
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    shapes = slide.shapes
    title = shapes.title
    title.text = "Gantt chart"
    left = Inches(1)
    top = Inches(1)
    pic = slide.shapes.add_picture("fig1.png", left, top,width = Inches(8),height = Inches(5.0))
    prs.save('./Final - Presentation.pptx')

    return True, "Created the Gantt chart"

details = {
    "gantt_data":{
            "Task":[1,2,3],
            "Start":['2009-01-01','2009-03-05','2009-02-20'],
            "Finish":['2009-02-28','2009-04-15','2009-07-15'],
            "Completion_pct":[50,25,75]
            }
        }
# print(create_gantt_chart(details))


req = { 
    "req":'create_chart',
    "details": {
        "create_gantt_chart":True,
        "gantt_data":{
            "Task":['a','b','c'],
            "Start":['2009-01-01','2009-03-05','2009-02-20'],
            "Finish":['2009-02-28','2009-04-15','2009-05-30'],
            "Completion_pct":[50,25,75]
        }
    }
}
