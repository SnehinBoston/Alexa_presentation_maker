from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm, Inches
import plotly.graph_objects as px

def create_stacked_bar_chart(details,base=None):
    print("I am here.")
    try:
        prs = Presentation('./Final - Presentation.pptx')
    except exc.PackageNotFoundError as err:
        prs = Presentation(base)
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    x = details['bar_stacked_data']['x']
    y = details['bar_stacked_data']['y']
    z = details['bar_stacked_data']['z'] 

    plot = px.Figure(data=[px.Bar(
        name = 'Data 1',
        x = x,
        y = y
    ),px.Bar(
        name = 'Data 2',
        x = x,
        y = z
    )])
    
    plot.update_layout(barmode='stack')
    
    plot.write_image("stacked_bar_chart.jpg")
    
    shapes = prs.slides[0].shapes
    picture = shapes.add_picture('stacked_bar_chart.jpg', Inches(0.1), Inches(0.1))

    prs.save('./Final - Presentation.pptx')

    return True, "Created the Stacked bar chart"


# data ={
#     "req":'create_stacked_bar_chart',
#     "details" : {
#         "bar_stacked_data":{
#             "x":['A','B','C'],
#             "y":[100,200],
#             "z":[125,225,325,425]
#         }
#     }
# }
# create_stacked_bar_chart(details)